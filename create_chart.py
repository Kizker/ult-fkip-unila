from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

def create_chart():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # blank slide with title

    # Define chart data
    chart_data = CategoryChartData()
    chart_data.categories = ['Ahli Materi', 'Ahli Media', 'Ahli Sistem']
    chart_data.add_series('Persentase Kelayakan (%)', (95.45, 93.33, 87.58))

    # Add chart to slide
    x, y, cx, cy = Inches(1.5), Inches(2), Inches(7), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    # Add title and format
    title = slide.shapes.title
    title.text = "Grafik Rekapitulasi Uji Validitas Ahli"
    
    chart.has_legend = True
    chart.legend.include_in_layout = False
    
    prs.save(r'C:\laragon\www\ult-fkip-unila\docs\skripsi\hasil_update\Grafik_Validator_Ahli.pptx')
    print("Grafik berhasil dibuat!")

create_chart()
