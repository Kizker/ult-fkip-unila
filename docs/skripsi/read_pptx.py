from pptx import Presentation
import sys

# Ensure utf-8 encoding for stdout
sys.stdout.reconfigure(encoding='utf-8')

try:
    prs = Presentation(r"c:\laragon\www\ult-fkip-unila\docs\skripsi\PPT Sempro Andri.pptx")
    print(f"Total slides: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                print(f"[TEXT] {shape.text.strip()}")
            elif shape.has_table:
                print("[TABLE]")
            elif str(shape.shape_type) == 'PICTURE (13)':
                print("[PICTURE]")
            else:
                print(f"[{shape.shape_type}]")
except Exception as e:
    print(f"Error: {e}")
