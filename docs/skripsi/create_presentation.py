from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Colors
DARK_BLUE = RGBColor(10, 34, 64)
WHITE = RGBColor(255, 255, 255)
LIGHT_BLUE = RGBColor(0, 112, 192)
GRAY = RGBColor(230, 230, 230)
DARK_GRAY = RGBColor(80, 80, 80)
ACCENT_ORANGE = RGBColor(237, 125, 49)
ACCENT_GREEN = RGBColor(112, 173, 71)

def add_title_slide(title, subtitle, author):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(22)
    p2.font.color.rgb = LIGHT_BLUE
    p2.alignment = PP_ALIGN.CENTER
    
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = author
    p3.font.size = Pt(18)
    p3.font.color.rgb = WHITE
    p3.alignment = PP_ALIGN.CENTER

    return slide

def add_chapter_divider(chapter_num, chapter_title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BLUE
    bg.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = chapter_num
    p.font.bold = True
    p.font.size = Pt(60)
    p.font.color.rgb = WHITE
    
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = chapter_title
    p2.font.bold = True
    p2.font.size = Pt(44)
    p2.font.color.rgb = DARK_BLUE
    
    return slide

def add_content_with_bullets(title, bullets, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title Ribbon
    ribbon = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.5), Inches(8), Inches(0.8))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = DARK_BLUE
    ribbon.line.fill.background()
    
    p = ribbon.text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    
    # Content Box
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, (bullet_title, bullet_desc) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet_title + ": "
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = DARK_BLUE
        p.level = 0
        
        run = p.add_run()
        run.text = bullet_desc
        run.font.bold = False
        run.font.size = Pt(22)
        run.font.color.rgb = DARK_GRAY
        
        # Add a bit of spacing between bullets
        if i < len(bullets) - 1:
            spacer = tf.add_paragraph()
            spacer.text = ""
            spacer.font.size = Pt(10)

    slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_mixed_layout_slide(title, main_stat, stat_label, bullets, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    p = txBox.text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = DARK_BLUE
    
    # Left Highlight Box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(3.5), Inches(5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = main_stat
    p1.font.bold = True
    p1.font.size = Pt(60)
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "\n" + stat_label
    p2.font.size = Pt(22)
    p2.alignment = PP_ALIGN.CENTER
    
    # Right Content Box
    content_box = slide.shapes.add_textbox(Inches(4.5), Inches(1.5), Inches(5), Inches(5))
    tf_c = content_box.text_frame
    tf_c.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            pc = tf_c.paragraphs[0]
        else:
            pc = tf_c.add_paragraph()
            # Spacing
            spacer = tf_c.add_paragraph()
            spacer.text = ""
            spacer.font.size = Pt(10)
            pc = tf_c.add_paragraph()
            
        pc.text = bullet
        pc.font.size = Pt(20)
        pc.font.color.rgb = DARK_GRAY
        pc.level = 0
        
    slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_image_text_slide(title, image_placeholder, bullets, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    p = txBox.text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = DARK_BLUE
    
    # Image Placeholder (Left)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY
    shape.line.fill.background()
    tf_s = shape.text_frame
    tf_s.word_wrap = True
    ps = tf_s.paragraphs[0]
    ps.text = f"<INSERT_IMAGE:\n{image_placeholder}>"
    ps.font.size = Pt(18)
    ps.font.color.rgb = DARK_GRAY
    ps.alignment = PP_ALIGN.CENTER
    
    # Content Box (Right)
    content_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.5), Inches(4.2), Inches(5))
    tf_c = content_box.text_frame
    tf_c.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            pc = tf_c.paragraphs[0]
        else:
            pc = tf_c.add_paragraph()
            spacer = tf_c.add_paragraph()
            spacer.text = ""
            spacer.font.size = Pt(10)
            pc = tf_c.add_paragraph()
            
        pc.text = bullet
        pc.font.size = Pt(20)
        pc.font.color.rgb = DARK_GRAY
        pc.level = 0
        
    slide.notes_slide.notes_text_frame.text = notes
    return slide

# ================= GENERATE SLIDES ================= #

# 1. Title
add_title_slide(
    "Pengembangan Website ULT FKIP Unila",
    "Digitalisasi Pelayanan Terpadu",
    "Andricha Dea Mitra\nNPM 2213025061"
)
prs.slides[-1].notes_slide.notes_text_frame.text = "Selamat pagi, Yang Terhormat Dewan Penguji, Bapak/Ibu Dosen Pembimbing, serta hadirin sekalian. Terima kasih atas waktu dan kesempatan yang diberikan kepada saya. Hari ini, saya Andricha Dea Mitra akan mempresentasikan hasil akhir penelitian skripsi saya yang berjudul 'Pengembangan Website Unit Layanan Terpadu (ULT) sebagai Upaya Digitalisasi Pelayanan di Fakultas Keguruan dan Ilmu Pendidikan Universitas Lampung'. Pada presentasi kali ini, saya akan menguraikan secara komprehensif mulai dari permasalahan dasar di lapangan, solusi arsitektur perangkat lunak yang saya bangun, hingga hasil akhir pengujian sistem ini oleh pakar dan pengguna langsung."

# 2. BAB I Pendahuluan
add_chapter_divider("BAB 1", "Pendahuluan")

# 3. Latar Belakang
add_mixed_layout_slide(
    "Latar Belakang & Urgensi Penelitian",
    "87,5%",
    "Mahasiswa mendukung urgensi digitalisasi layanan terpadu berbasis website.",
    [
        "Permasalahan Utama: Pelayanan dokumen secara manual di ULT memicu antrean panjang, ketidakjelasan status dokumen, dan tingginya risiko dokumen fisik yang tercecer atau hilang.",
        "Kondisi Saat Ini: Survei awal menunjukkan bahwa tingkat kepuasan layanan konvensional berada di bawah 50%.",
        "Solusi yang Ditawarkan: Transformasi ke sistem digital berbasis website yang transparan, mudah diakses, dan mampu melacak setiap pergerakan dokumen."
    ],
    "Kita mulai dari latar belakang masalah. Penelitian ini lahir dari realitas lapangan di mana pelayanan administrasi mahasiswa di ULT saat ini masih mengandalkan proses manual. Mahasiswa harus datang langsung untuk menyerahkan berkas fisik. Proses manual ini memicu penumpukan antrean, hilangnya berkas, dan mahasiswa tidak bisa memantau sudah sejauh mana surat mereka diproses. Dari survei pendahuluan yang saya lakukan, kepuasan terhadap sistem lama sangat rendah, dan 87,5 persen mahasiswa dengan tegas menyatakan bahwa digitalisasi adalah sebuah urgensi. Oleh karena itu, penelitian ini hadir untuk menawarkan solusi berupa sebuah sistem website yang transparan, cepat, dan terpusat."
)

# 4. Rumusan Masalah
add_content_with_bullets(
    "Rumusan Masalah",
    [
        ("1. Proses Pengembangan (ADDIE)", "Bagaimana proses pengembangan website ULT FKIP Universitas Lampung menggunakan pendekatan R&D dengan model ADDIE?"),
        ("2. Tingkat Kelayakan/Validitas", "Bagaimana tingkat kelayakan atau validitas website ULT FKIP Universitas Lampung berdasarkan hasil validasi ahli materi, ahli media dan ahli sistem?"),
        ("3. Tingkat Kepraktisan", "Bagaimana tingkat kepraktisan penggunaan website ULT FKIP Universitas Lampung berdasarkan hasil angket uji kepraktisan?")
    ],
    "Dari latar belakang yang telah dipaparkan, saya menarik tiga rumusan masalah utama yang sama persis dengan dokumen skripsi. Pertama, menanyakan bagaimana proses pengembangannya menggunakan metode R&D model ADDIE. Kedua, menanyakan tingkat kelayakan dan validitasnya di mata para ahli. Dan yang ketiga, menanyakan bagaimana tingkat kepraktisan sistem ini saat digunakan langsung oleh pengguna melalui uji kepraktisan."
)

# 5. Tujuan Penelitian
add_content_with_bullets(
    "Tujuan Penelitian",
    [
        ("1. Mengembangkan Website (ADDIE)", "Untuk mengembangkan website ULT FKIP Universitas Lampung menggunakan pendekatan R&D dengan model ADDIE."),
        ("2. Menilai Kelayakan/Validitas", "Untuk menilai tingkat kelayakan/validitas website ULT FKIP Universitas Lampung berdasarkan hasil validasi ahli materi, ahli media dan ahli sistem."),
        ("3. Mengukur Tingkat Kepraktisan", "Untuk mengukur tingkat kepraktisan penggunaan website ULT FKIP Universitas Lampung berdasarkan hasil angket uji kepraktisan.")
    ],
    "Sejalan dengan rumusan masalah tersebut, maka tujuan dari penelitian ini juga terbagi menjadi tiga sasaran utama yang linier. Tujuan pertama adalah untuk mengembangkan website tersebut menggunakan model ADDIE. Tujuan kedua adalah untuk menilai kelayakan dan validitasnya dari kacamata para ahli. Dan tujuan ketiga adalah untuk mengukur seberapa praktis website ini saat digunakan melalui penyebaran angket kepraktisan."
)

# 6. BAB II Tinjauan Pustaka
add_chapter_divider("BAB 2", "Tinjauan Pustaka")

# 7. Tinjauan Pustaka Content
add_content_with_bullets(
    "Landasan Teori Utama",
    [
        ("Metodologi Penelitian", "Menggunakan pendekatan Research & Development (R&D) dengan model ADDIE (Analisis, Desain, Pengembangan, Implementasi, dan Evaluasi)."),
        ("Teknologi & Keamanan", "Sistem dikembangkan menggunakan Framework Laravel dengan pengamanan akses berlapis (Role-Based Access Control) untuk melindungi privasi data mahasiswa."),
        ("Pengukuran Kelayakan", "Validitas pakar diukur menggunakan rumus koefisien Aiken's V, dan kepraktisan pengguna diukur melalui kuesioner berskala standar.")
    ],
    "Dalam mengembangkan sistem ini, saya berpijak pada tiga pilar landasan teori utama. Pertama, dari segi metodologi, saya menggunakan model Research & Development (R&D) dengan kerangka ADDIE karena model ini sangat sistematis dan cocok untuk rekayasa perangkat lunak. Kedua, dari segi teknologi, saya memilih Framework Laravel karena keandalannya dalam menangani struktur data besar dan kemampuannya dalam menerapkan sistem keamanan tingkat tinggi seperti Role-Based Access Control. Terakhir, untuk menjamin kualitas produk, saya menggunakan formulasi statistik Aiken's V untuk mengukur tingkat validitas ahli, sehingga produk yang dihasilkan benar-benar teruji secara akademis dan teknis."
)

# 8. Penelitian Terdahulu
add_content_with_bullets(
    "Penelitian Terdahulu yang Relevan",
    [
        ("Fokus Terhadap Layanan Digital Akademik", "Penelitian terdahulu banyak berfokus pada digitalisasi layanan, namun sebatas pada perpindahan dari form kertas ke platform konvensional pihak ketiga (misal: Google Form)."),
        ("Sistem Informasi Monolitik Lama", "Beberapa pengembangan sebelumnya berhasil membangun sistem internal, namun seringkali mengabaikan aspek isolasi privasi antar pengguna."),
        ("Keterbatasan Automasi Dokumen", "Mayoritas studi terdahulu memiliki output di mana dokumen cetak masih harus dirakit dan diformat secara manual oleh staf administrasi.")
    ],
    "Beralih ke kajian penelitian terdahulu. Dari studi literatur yang saya lakukan, banyak penelitian sebelumnya yang telah mengangkat tema digitalisasi pelayanan akademik. Namun, kebanyakan dari sistem terdahulu masih terjebak pada digitalisasi semu—seperti sekadar memindahkan form cetak menjadi form online biasa tanpa terintegrasi. Ada pula yang sudah membangun sistem utuh, namun sering kali mengabaikan isolasi keamanan privasi data mahasiswanya. Selain itu, kelemahan terbesar sistem pendahulu adalah ketiadaan automasi dokumen, sehingga pada akhirnya staf masih harus menyalin dan merakit isian dokumen kelulusan mahasiswa secara manual satu per satu."
)

# 9. Gap Penelitian
add_content_with_bullets(
    "Gap Penelitian (Mengapa Penelitian Ini Dilakukan)",
    [
        ("Inovasi Perakitan Dokumen Otomatis", "Sistem ini membawa pembaruan mutlak berupa mesin 'Document Assembly' yang merakit teks mahasiswa ke dalam template MS Word secara presisi dan instan."),
        ("Pengamanan Privasi Ketat (Anti-IDOR)", "Berbeda dengan sistem pendahulu, website ULT ini menerapkan proteksi Anti-IDOR agar mahasiswa tidak bisa mengintip berkas pribadi milik orang lain."),
        ("Transparansi Siklus (Audit Trail)", "Menjawab keluhan mahasiswa, sistem ini memungkinkan pelacakan jejak pergerakan surat secara real-time dari awal pengajuan hingga surat selesai.")
    ],
    "Berangkat dari berbagai keterbatasan penelitian terdahulu tersebut, di sinilah letak 'Gap Penelitian' yang mendasari mengapa penelitian saya ini wajib untuk dilakukan. Penelitian ini hadir untuk menutup kelemahan-kelemahan masa lalu dengan menawarkan tiga pembaruan inovatif. Pembaruan pertama, tidak ada lagi perakitan dokumen secara manual; sistem ini secara pintar merakit form isian mahasiswa ke template surat kampus secara instan. Kedua, sistem ini dijaga dengan perlindungan Anti-IDOR yang mengunci keamanan data sehingga privasi tak tertembus. Dan ketiga, penerapan fitur lacak dokumen transparan yang memungkinkan mahasiswa memantau pergerakan surat mereka seperti sedang melacak paket kiriman. Pembaruan inilah yang membuat sistem ULT ini lebih superior dibanding pendahulunya."
)

# 6. BAB III Metodologi
add_chapter_divider("BAB 3", "Metodologi Penelitian")

# 7. Metodologi Content
add_mixed_layout_slide(
    "Prosedur & Metodologi",
    "ADDIE",
    "Analysis, Design, Development, Implementation, Evaluation",
    [
        "Pendekatan Penelitian: Mengadopsi metode Research & Development (R&D) secara penuh dari awal hingga produk siap pakai.",
        "Populasi & Sampel: Dipilih melalui teknik purposive sampling untuk mendapatkan penilai yang benar-benar kredibel.",
        "Subjek Pengujian: Melibatkan 9 Pakar Ahli (untuk validasi teknis dan materi) serta 18 Pengguna Akhir (untuk menguji kenyamanan dan kepraktisan penggunaan)."
    ],
    "Memasuki bab metodologi. Pengembangan website ini dijalankan secara berurutan melewati 5 tahapan ADDIE: mulai dari Analisis Kebutuhan, Desain Sistem, Pengembangan Kode, Implementasi di lapangan, hingga Evaluasi akhir. Untuk memastikan sistem ini valid dan tidak sekadar jadi, saya melibatkan total 27 responden yang dipilih secara spesifik (purposive sampling). 9 orang di antaranya adalah pakar akademisi ahli di bidang materi, media, dan sistem yang bertugas membedah jeroan teknis website. Sedangkan 18 orang lainnya adalah mahasiswa, staf ULT, dan admin prodi yang bertindak sebagai penguji langsung di lapangan untuk menilai apakah website ini benar-benar mempermudah pekerjaan mereka atau tidak."
)

# 8. BAB IV Hasil dan Pembahasan
add_chapter_divider("BAB 4", "Hasil dan Pembahasan")

# 9. Tahap Analisis
add_image_text_slide(
    "Hasil: Tahap Analisis (Analysis)",
    "Ilustrasi/Data\nKebutuhan Sistem",
    [
        "Identifikasi Masalah: Menemukan kelemahan fatal pada sistem manual, yakni ketiadaan rekam jejak digital yang transparan bagi mahasiswa.",
        "Pengumpulan Data: Observasi dan wawancara langsung dengan staf ULT mengonfirmasi beban kerja administratif yang terlalu berat.",
        "Hasil Kebutuhan: Mahasiswa dan staf sama-sama membutuhkan platform yang bisa diakses 24 jam dari mana saja secara online."
    ],
    "Masuk ke bagian utama, yakni Hasil Penelitian yang saya jabarkan mengikuti alur 5 tahap ADDIE. Tahap pertama adalah Analisis. Pada tahap ini, saya mendalami letak permasalahan sistem manual. Ditemukan bahwa staf ULT sering kewalahan menyortir ratusan kertas, sementara mahasiswa kebingungan mencari tahu status dokumen mereka. Berdasarkan pengumpulan data dan wawancara, ditarik kesimpulan bahwa kebutuhan utamanya adalah sebuah platform digital online yang mampu menampung pengajuan 24 jam sehari, serta mampu melacak keberadaan dokumen secara otomatis."
)

# 10. Tahap Desain
add_image_text_slide(
    "Hasil: Tahap Desain (Design)",
    "Diagram Arsitektur Sistem\natau Flowchart Dokumen",
    [
        "Arsitektur Sistem: Sistem dibangun dengan rancangan pembagian portal akses yang ketat (Public, Student, Admin, Staff, dan Signer).",
        "Skema Database: Merancang struktur database dengan formulir dinamis, sehingga staf ULT kelak dapat menambah layanan baru tanpa perlu membongkar kode program.",
        "Alur Kerja (Workflow): Menetapkan alur perjalanan dokumen secara terstruktur mulai dari draf, ditinjau, hingga disetujui."
    ],
    "Tahap kedua adalah Desain, yaitu menerjemahkan kebutuhan tadi ke dalam sebuah cetak biru arsitektur. Saya merancang sistem ini agar membedakan akses penggunanya secara tegas ke dalam 5 portal (untuk Publik, Mahasiswa, Admin, Staf, dan Penandatangan). Keunggulan dari desain ini adalah skema database yang dinamis. Artinya, jika suatu hari kampus ingin menambah jenis layanan surat baru, staf ULT dapat langsung membuatnya dari dalam aplikasi tanpa perlu memanggil programmer untuk merombak kode database. Saya juga mendesain alur kerja agar setiap perpindahan dokumen dicatat secara digital."
)

# 11. Tahap Pengembangan
add_mixed_layout_slide(
    "Hasil: Tahap Pengembangan (Development)",
    "91,95%",
    "Tingkat Validasi Pakar Ahli (Sangat Valid)",
    [
        "Realisasi Sistem: Desain diwujudkan ke dalam kode pemrograman web menggunakan teknologi terkini (Laravel 12, PHP 8.4, dan TailwindCSS).",
        "Pengujian Ketat: Setelah aplikasi berdiri, sistem dibedah dan diuji secara komprehensif oleh 9 Pakar Ahli independen.",
        "Skor Kelayakan: Meraih persentase total 91,95% dengan tingkat keandalan statistik (Aiken's V) sebesar 0,90 yang berarti 'Sangat Valid'."
    ],
    "Tahap ketiga adalah Pengembangan. Di sinilah desain dieksekusi menjadi baris-baris kode pemrograman fungsional menggunakan framework mutakhir Laravel 12. Setelah purwarupa website ini jadi, saya tidak langsung melepasnya ke mahasiswa, melainkan menyerahkannya kepada 9 pakar ahli untuk diverifikasi secara teknis. Para pakar ini menguji keamanan, tampilan, serta akurasi fungsionalnya. Hasil pengujian sungguh memuaskan, mencatatkan skor 91,95 persen. Angka ini secara statistik mengukuhkan bahwa perangkat lunak yang saya bangun sangat solid, minim celah bahaya, dan dinyatakan 'Sangat Valid' untuk digunakan."
)

# 12. Tahap Implementasi
add_mixed_layout_slide(
    "Hasil: Tahap Implementasi (Implementation)",
    "92,13%",
    "Tingkat Kepraktisan Pengguna (Sangat Praktis)",
    [
        "Uji Coba Lapangan: Mengundang 18 pengguna akhir (mahasiswa dan staf) untuk mencoba langsung proses pengajuan dan pemrosesan surat.",
        "Tingkat Kenyamanan: Para pengguna menilai sistem ini sangat intuitif dan mudah dipahami, tercermin dari perolehan skor tinggi 92,13%.",
        "Penyempurnaan Langsung: Umpan balik dari staf langsung diterapkan ke dalam sistem, seperti penambahan fitur nomor surat yang terisi otomatis."
    ],
    "Setelah lolos uji ahli, sistem memasuki tahap keempat: Implementasi. Sistem diterjunkan ke kondisi lapangan sesungguhnya dengan melibatkan 18 responden yang bertindak sebagai pengguna harian (mahasiswa dan staf). Pengujian ini berfokus pada kenyamanan dan kemudahan (usabilitas). Hasil kuesioner menunjukkan skor kepraktisan 92,13 persen (Sangat Praktis). Di tahap ini, saya juga secara langsung menampung keluhan staf dan seketika memperbarui sistem—salah satu contohnya adalah menambahkan algoritma agar sistem dapat memberikan nomor surat secara otomatis, sehingga meringankan beban kerja operasional staf ULT."
)

# 13. Tahap Evaluasi
add_content_with_bullets(
    "Hasil: Tahap Evaluasi (Evaluation)",
    [
        ("Evaluasi Berkelanjutan", "Proses perbaikan selalu dilakukan di setiap celah yang ditemukan selama tahapan sebelumnya, khususnya pada perkuatan sistem keamanan web."),
        ("Evaluasi Akhir", "Analisis final menyimpulkan bahwa seluruh fitur website beroperasi dengan sempurna dan tanpa cacat yang mengganggu fungsionalitas."),
        ("Kesuksesan Target", "Sistem secara definitif terbukti memecahkan masalah birokrasi layanan manual yang selama ini dikeluhkan oleh mahasiswa.")
    ],
    "Tahapan kelima atau terakhir dari ADDIE adalah Evaluasi. Proses evaluasi ini sebenarnya sudah berjalan sejak awal pengembangan untuk langsung menambal kelemahan yang ditemukan (evaluasi formatif). Pada evaluasi akhir (sumatif), kami menarik kesimpulan bulat bahwa keseluruhan rancang bangun produk website ULT FKIP Unila ini berhasil beroperasi dengan sangat baik. Sistem web ini secara definitif terbukti efektif dalam memangkas jalur birokrasi, mengamankan privasi dokumen mahasiswa, dan menjawab langsung permasalahan inti yang memicu penelitian ini."
)

# 14. Pembahasan: Keunggulan
add_content_with_bullets(
    "Pembahasan: Keunggulan Sistem",
    [
        ("Isolasi Keamanan Ketat", "Dokumen mahasiswa dikunci di dalam brankas digital (private folder) yang dilindungi sistem pencegah peretasan (Anti-IDOR)."),
        ("Perakitan Dokumen Otomatis", "Sistem mampu merakit tulisan isian form mahasiswa ke dalam template surat cetak secara instan, tanpa membuat susunan teks menjadi berantakan."),
        ("Transparansi Lacak Dokumen", "Sama seperti melacak paket pengiriman online, mahasiswa dapat melihat secara real-time sudah di tahap mana surat mereka diproses (Auditable Timeline).")
    ],
    "Beralih ke bagian Pembahasan. Di sini saya ingin menekankan 3 keunggulan mutlak yang membuat sistem ini istimewa. Pertama, keamanan ekstra ketat; saya menyematkan perlindungan Anti-IDOR agar tidak ada pihak luar yang bisa mengintip atau mencuri file dokumen pribadi milik mahasiswa lain. Kedua, fitur 'Mesin Perakitan Dokumen'; di mana sistem secara pintar menyatukan isian form mahasiswa ke dalam format resmi surat kampus Microsoft Word tanpa merusak posisi huruf, margin, atau spasinya sama sekali. Ketiga, transparansi penuh; mahasiswa kini bisa melacak dokumen mereka seakan sedang melacak posisi paket kurir secara real-time."
)

# 15. Pembahasan: Kendala & Solusi
add_content_with_bullets(
    "Pembahasan: Kendala & Solusi Teknis",
    [
        ("Teks Berantakan Saat Dicetak", "Teks dari website awalnya selalu merusak format file Word. Solusi: Saya merancang penerjemah khusus (HtmlToOpenXMLParser) agar formatnya tersusun rapi secara otomatis."),
        ("Ancaman Virus/Kode Jahat", "Celah pengisian form rentan disusupi script berbahaya. Solusi: Sistem dibekali alat penyaring canggih (HTMLSanitizer) untuk menetralisir ancaman keamanan."),
        ("Belum Tersambung ke Universitas", "Sistem saat ini masih berdiri sendiri di level Fakultas. Solusi: Pondasi aplikasi telah dipersiapkan agar suatu hari mudah digabungkan dengan sistem login utama Universitas (SSO).")
    ],
    "Tentu, ada tantangan teknis yang berat selama pengerjaan. Kendala terbesar adalah ketika mengonversi form isian website ke dokumen Microsoft Word, format visual aslinya selalu hancur. Untuk menanganinya, saya mendesain algoritma penerjemah khusus (HtmlToOpenXMLParser) yang memaksa format website tunduk pada format cetakan Word secara rapi. Selain itu, ada kerentanan keamanan di mana orang iseng bisa memasukkan kode peretasan melalui form; ini saya redam dengan mengaktifkan penyaring keamanan web mutakhir. Terakhir, karena sistem ini belum terhubung langsung ke server pusat Unila, saya sudah membekali kodenya dengan fondasi arsitektur siap-pakai agar sewaktu-waktu bisa segera dihubungkan ke sistem Universitas (SSO)."
)

# 16. BAB V Kesimpulan & Saran
add_chapter_divider("BAB 5", "Kesimpulan dan Saran")

# 17. Kesimpulan
add_content_with_bullets(
    "Kesimpulan Penelitian",
    [
        ("1. Proses Pengembangan (ADDIE)", "Website ULT berhasil dikembangkan dengan menerapkan model pengembangan ADDIE secara sistematis, menghasilkan aplikasi web monolitik yang aman dan terpusat."),
        ("2. Tingkat Kelayakan/Validitas", "Uji kelayakan oleh 9 validator ahli menunjukkan rata-rata keseluruhan 91,95% yang dikategorikan 'Sangat Valid'."),
        ("3. Tingkat Kepraktisan", "Uji coba terbatas pada responden civitas akademika menghasilkan persentase kepraktisan 92,13% yang menempatkan sistem ke dalam kriteria 'Sangat Praktis'.")
    ],
    "Sebagai penutup presentasi ini, Kesimpulan dari penelitian ini menjawab ketiga rumusan masalah secara lugas yang sejalan dengan tujuan penelitian. Pertama, untuk proses pengembangannya, website ULT ini telah berhasil dikembangkan dengan mantap dan sistematis menggunakan model ADDIE menjadi sebuah aplikasi yang aman. Kedua, dari aspek kelayakan dan validitas, pengujian dari para ahli pakar telah membuktikan bahwa sistem ini meraih persentase 91,95 persen dengan predikat Sangat Valid. Dan ketiga, dari aspek kepraktisan di lapangan, sistem ini telah dibuktikan secara empiris meraih tingkat kepraktisan 92,13 persen dengan predikat Sangat Praktis oleh civitas akademika penggunanya."
)

# 18. Saran
add_content_with_bullets(
    "Saran Pengembangan Selanjutnya",
    [
        ("Ekspansi Layanan Universitas", "Direkomendasikan agar sistem ini terus dikembangkan agar kelak dapat digunakan langsung menggunakan akun mahasiswa terpusat Unila (Single Sign-On)."),
        ("Penambahan Modul Staf", "Menyarankan adanya penambahan portal fitur khusus yang dapat melayani administrasi persuratan internal untuk kebutuhan dosen dan pegawai staf.")
    ],
    "Berangkat dari hasil pengembangan ini, saya menyarankan dua poin pengembangan di masa depan. Pertama, integrasi lanjutan ke sistem pusat, agar mahasiswa dapat login menggunakan akun pusat kampus (SSO Unila). Kedua, pengembangan agar platform ini tidak hanya melayani permohonan administrasi mahasiswa, tetapi juga mengakomodasi pelayanan administrasi persuratan bagi bapak ibu dosen dan pegawai fakultas. Demikianlah presentasi hasil penelitian skripsi saya. Terima kasih atas perhatian Dewan Penguji sekalian. Waktu dan tempat saya kembalikan kepada moderator."
)

prs.save('c:\\laragon\\www\\ult-fkip-unila\\docs\\skripsi\\Presentasi_Seminar_Hasil_ULT_Andricha.pptx')
print("Presentation regenerated with detailed contents successfully!")
